"""
Parseurs de documents — extraction du texte depuis différents formats.

Chaque parseur retourne une liste de tuples (page_number, text).
Le numéro de page est utilisé pour les citations dans la réponse finale.

Pour les formats sans pagination native (docx, txt), on utilise des
"blocs logiques" : paragraphes ou sections.
"""
from __future__ import annotations

import io
import os
from pathlib import Path
from typing import List, Tuple

# Seuil en dessous duquel on considère qu'un PDF est scanné (pas de texte extractible)
OCR_FALLBACK_THRESHOLD_CHARS = 100


def parse_document(file_path: str, filename: str) -> List[Tuple[int, str]]:
    """
    Router principal : détecte le format et dispatche vers le bon parseur.

    Args:
        file_path: chemin absolu du fichier sur disque.
        filename: nom du fichier original (pour détecter l'extension).

    Returns:
        Liste de (page_number, text). Les pages vides sont filtrées.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        pages = _parse_pdf(file_path)
        # Si le PDF contient peu de texte, on tente l'OCR
        total_chars = sum(len(t) for _, t in pages)
        if total_chars < OCR_FALLBACK_THRESHOLD_CHARS:
            print(f"[ingestion] PDF peu de texte ({total_chars} chars) → fallback OCR")
            pages = _parse_pdf_ocr(file_path)
        return pages

    if ext == ".docx":
        return _parse_docx(file_path)

    if ext == ".pptx":
        return _parse_pptx(file_path)

    if ext in (".txt", ".md"):
        return _parse_text(file_path)

    if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"):
        return _parse_image_ocr(file_path)

    raise ValueError(
        f"Format non supporté : {ext}. "
        "Formats acceptés : PDF, DOCX, PPTX, TXT, MD, PNG, JPG, JPEG, WEBP, BMP, TIFF."
    )


# ============================================
# PDF — texte natif
# ============================================
def _parse_pdf(file_path: str) -> List[Tuple[int, str]]:
    """Extraction texte d'un PDF via pdfplumber (meilleur que pypdf pour la mise en page)."""
    import pdfplumber

    pages: List[Tuple[int, str]] = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                pages.append((i, text))
    return pages


# ============================================
# PDF — fallback OCR (pour les cours scannés)
# ============================================
def _parse_pdf_ocr(file_path: str) -> List[Tuple[int, str]]:
    """Convertit chaque page PDF en image puis applique Tesseract OCR."""
    from pdf2image import convert_from_path
    import pytesseract

    # Langues : arabe + français + anglais (à installer côté OS)
    lang = "ara+fra+eng"

    pages: List[Tuple[int, str]] = []
    images = convert_from_path(file_path, dpi=200)
    for i, img in enumerate(images, start=1):
        try:
            text = pytesseract.image_to_string(img, lang=lang)
        except pytesseract.TesseractError:
            # Si les packs de langue ne sont pas installés, on retombe sur anglais
            text = pytesseract.image_to_string(img)
        text = (text or "").strip()
        if text:
            pages.append((i, text))
    return pages


# ============================================
# Image seule — OCR direct
# ============================================
def _parse_image_ocr(file_path: str) -> List[Tuple[int, str]]:
    from PIL import Image
    import pytesseract

    lang = "ara+fra+eng"
    img = Image.open(file_path)
    try:
        text = pytesseract.image_to_string(img, lang=lang)
    except pytesseract.TesseractError:
        text = pytesseract.image_to_string(img)
    text = (text or "").strip()
    return [(1, text)] if text else []


# ============================================
# DOCX — Word
# ============================================
def _parse_docx(file_path: str) -> List[Tuple[int, str]]:
    """
    Word n'a pas de notion de page native en XML. On simule la pagination en
    regroupant ~30 paragraphes par "page logique" pour avoir des citations
    raisonnables.
    """
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Également extraire les tableaux
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    if not paragraphs:
        return []

    # Grouper par "pages logiques" de ~30 paragraphes
    PARA_PER_PAGE = 30
    pages: List[Tuple[int, str]] = []
    for page_num, start in enumerate(range(0, len(paragraphs), PARA_PER_PAGE), start=1):
        chunk = "\n".join(paragraphs[start : start + PARA_PER_PAGE])
        if chunk:
            pages.append((page_num, chunk))
    return pages


# ============================================
# PPTX — PowerPoint
# ============================================
def _parse_pptx(file_path: str) -> List[Tuple[int, str]]:
    """Une slide = une page logique. Idéal pour les citations."""
    from pptx import Presentation

    prs = Presentation(file_path)
    pages: List[Tuple[int, str]] = []

    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            # Text boxes
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        texts.append(txt)
            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_txt = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_txt:
                        texts.append(row_txt)

        # Speaker notes (souvent précieuses pour un cours)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Notes du présentateur] {notes}")

        if texts:
            pages.append((i, "\n".join(texts)))

    return pages


# ============================================
# TXT / MD
# ============================================
def _parse_text(file_path: str) -> List[Tuple[int, str]]:
    """Fichier texte : on crée des pages logiques de ~2000 caractères."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read().strip()

    if not content:
        return []

    CHARS_PER_PAGE = 2000
    pages: List[Tuple[int, str]] = []
    for page_num, start in enumerate(range(0, len(content), CHARS_PER_PAGE), start=1):
        chunk = content[start : start + CHARS_PER_PAGE]
        pages.append((page_num, chunk))
    return pages
